#!/usr/bin/env python3
"""
Generate the Zava Learning SRE-Agent lab architecture deck (Azure blue/white theme).

Produces docs/Zava-Learning-Architecture.pptx:
  * platform + network architecture
  * the incident lifecycle (fault -> alert -> PagerDuty -> SRE Agent -> IaC PR)
  * how the SRE Agent is wired (custom agent, skills, connectors, tools) + runbook/artifacts
  * the 7 parallel fault "lanes" map
  * per scenario: a good-state vs bad-state slide AND a response slide
    (how the fault is created, how the agent is triggered, what it uses, what it produces)

Run:  python docs/build_architecture_deck.py
"""
from pathlib import Path
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ── Azure blue/white palette ────────────────────────────────────────
AZURE      = RGBColor(0x00, 0x78, 0xD4)   # primary Azure blue
NAVY       = RGBColor(0x00, 0x20, 0x50)   # deep header navy
DARKBLUE   = RGBColor(0x24, 0x3A, 0x5E)
MIDBLUE    = RGBColor(0x10, 0x5E, 0xA6)
CYAN       = RGBColor(0x50, 0xE6, 0xFF)   # accent
LIGHT1     = RGBColor(0xEF, 0xF6, 0xFC)   # very light blue fill
LIGHT2     = RGBColor(0xC7, 0xE0, 0xF4)   # light blue fill
LIGHT3     = RGBColor(0xDE, 0xEC, 0xF9)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x23, 0x23, 0x23)
GRAY       = RGBColor(0x60, 0x60, 0x60)
GREEN      = RGBColor(0x10, 0x7C, 0x10)
GREENFILL  = RGBColor(0xDF, 0xF6, 0xDD)
RED        = RGBColor(0xC5, 0x29, 0x2C)
REDFILL    = RGBColor(0xFD, 0xE7, 0xE9)
AMBER      = RGBColor(0xC1, 0x6C, 0x00)

EMU_IN = 914400


def _in(v):
    return Emu(int(v * EMU_IN))


def add_box(slide, x, y, w, h, text="", fill=WHITE, line=AZURE, line_w=1.25,
            size=12, color=INK, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sub=None, sub_size=9,
            sub_color=None):
    sp = slide.shapes.add_shape(shape, _in(x), _in(y), _in(w), _in(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = _in(0.05)
    tf.margin_right = _in(0.05)
    tf.margin_top = _in(0.02)
    tf.margin_bottom = _in(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = "Segoe UI"
    r.font.color.rgb = color
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = align
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(sub_size)
        r2.font.name = "Segoe UI"
        r2.font.color.rgb = sub_color or GRAY
    return sp


def add_text(slide, x, y, w, h, text, size=12, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, name="Segoe UI"):
    tb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = name
        r.font.color.rgb = color
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=AZURE, width=1.75, dashed=False,
              head=False, tail=True):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _in(x1), _in(y1), _in(x2), _in(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    if dashed:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)
    if tail:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if head:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, _in(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = _in(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.name = "Segoe UI"; r.font.color.rgb = WHITE
    # accent stripe
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, _in(0.9), prs.slide_width, _in(0.06))
    acc.fill.solid(); acc.fill.fore_color.rgb = CYAN; acc.line.fill.background(); acc.shadow.inherit = False
    if subtitle:
        add_text(slide, 0.5, 1.02, 12.4, 0.4, subtitle, size=13, color=MIDBLUE, bold=True)


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    return s


def legend_chip(slide, x, y, color, label, w=2.0):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, _in(x), _in(y + 0.03), _in(0.16), _in(0.16))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background(); c.shadow.inherit = False
    add_text(slide, x + 0.24, y - 0.04, w, 0.3, label, size=10, color=INK)


def panel(slide, x, y, w, h, title, lines, num=None, title_fill=NAVY, title_color=CYAN,
          body_fill=LIGHT1, body_color=INK, line=AZURE, body_size=9):
    th = 0.42
    head = f"{num}  {title}" if num else title
    add_box(slide, x, y, w, th, head, fill=title_fill, line=None, size=11.5, bold=True,
            color=title_color, align=PP_ALIGN.LEFT)
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _in(x), _in(y + th), _in(w), _in(h - th))
    body.fill.solid(); body.fill.fore_color.rgb = body_fill
    body.line.color.rgb = line; body.line.width = Pt(1.0); body.shadow.inherit = False
    add_text(slide, x + 0.12, y + th + 0.08, w - 0.24, h - th - 0.16, "\n".join(lines),
             size=body_size, color=body_color)
    return body


prs = Presentation()
prs.slide_width = _in(13.333)
prs.slide_height = _in(7.5)


# ════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════════════════════════════
s = blank()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background(); band.shadow.inherit = False
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, _in(4.55), prs.slide_width, _in(0.09))
strip.fill.solid(); strip.fill.fore_color.rgb = CYAN; strip.line.fill.background(); strip.shadow.inherit = False
add_text(s, 0.9, 1.7, 11.5, 1.2, "Zava Learning", size=54, color=WHITE, bold=True)
add_text(s, 0.9, 2.85, 11.5, 1.0, "Azure SRE Agent Training Lab — Architecture & Fault Scenarios",
         size=24, color=CYAN, bold=True)
add_text(s, 0.92, 4.7, 11.5, 1.6,
         "A realistic, DB-backed online-learning platform with TWO SRE-Agent stories on one Azure\n"
         "environment:  (1) 7 parallel fault \"lanes\" + a back-office reporting worker that each plant a\n"
         "real, symptom-only fault for the agent to autonomously detect, diagnose, and remediate; and\n"
         "(2) weekly READ-ONLY governance audits (NSG · RBAC · cost) delivered as branded PowerPoint.",
         size=14, color=LIGHT2)
add_text(s, 0.92, 6.75, 12, 0.4, "Blue/White Azure theme  •  App Gateway · Container Apps · PostgreSQL · Key Vault · SRE Agent · PagerDuty · Scheduled Tasks",
         size=11, color=LIGHT2, italic=True)


# ════════════════════════════════════════════════════════════════════
# Slide 2 — Business context
# ════════════════════════════════════════════════════════════════════
s = blank()
header(s, "The Platform — Zava Learning", "A McGraw-Hill-style online learning product (students take timed quizzes)")
cards = [
    ("Learner Portal", "Public web app students browse during midterms.", AZURE),
    ("Course & Assessment APIs", "Serve course catalog and quiz content.", MIDBLUE),
    ("Quiz Service (DB-backed)", "Reads a 500k-row question bank from PostgreSQL.", AZURE),
    ("Gradebook API", "Writes & reads student submissions.", MIDBLUE),
]
x = 0.6
for title, desc, col in cards:
    add_box(s, x, 1.5, 2.9, 1.7, title, fill=LIGHT1, line=col, size=14, bold=True, color=NAVY,
            sub="\n" + desc, sub_size=10)
    x += 3.05
add_text(s, 0.6, 3.6, 12.1, 0.5,
         "Every student request flows through one public Application Gateway into Azure Container Apps,\n"
         "which read and write a managed PostgreSQL database. Secrets live in Key Vault.",
         size=13, color=INK)
# why
add_box(s, 0.6, 4.45, 12.1, 2.6,
        "Why this lab is \"almost real world\"",
        fill=NAVY, line=None, size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT)
add_text(s, 0.9, 5.05, 11.6, 1.95,
         "•  Faults are driven by REAL inputs — IaC parameters, real source → real image, real DB state,\n"
         "    or real Key Vault secrets — never env-var toggles or code-injection.\n"
         "•  Alerts and student symptoms are SYMPTOM-ONLY; they never name the cause.\n"
         "•  Diagnosing the root cause and proposing the durable IaC fix is the SRE Agent's job.\n"
         "•  7 quiz lanes run in parallel on ONE environment, plus a reporting-worker VM scenario.\n"
         "•  Beyond incidents, the same agent runs weekly READ-ONLY governance audits (NSG · RBAC ·\n"
         "    cost) delivered as branded PowerPoint — proactive posture, not just reactive response.",
         size=12.5, color=WHITE)


# ════════════════════════════════════════════════════════════════════
# Slide 3 — End-to-end architecture
# ════════════════════════════════════════════════════════════════════
s = blank()
header(s, "End-to-End Architecture", "How every component connects on a single Azure environment")

# Students
add_box(s, 0.45, 2.9, 1.5, 1.0, "Students", fill=AZURE, line=NAVY, size=14, bold=True, color=WHITE,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
# App Gateway
add_box(s, 2.45, 2.5, 1.9, 1.8, "Application\nGateway", fill=AZURE, line=NAVY, size=14, bold=True,
        color=WHITE, sub=":80 portal\n:8081–:8087 lanes", sub_size=9, sub_color=LIGHT2)
add_arrow(s, 1.95, 3.4, 2.45, 3.4, color=NAVY, width=2.25)

# VNet container
vnet = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _in(4.7), _in(1.35), _in(5.7), _in(5.65))
vnet.fill.solid(); vnet.fill.fore_color.rgb = LIGHT1; vnet.line.color.rgb = AZURE; vnet.line.width = Pt(1.5)
vnet.line.dash_style = None; vnet.shadow.inherit = False
add_text(s, 4.85, 1.42, 5.4, 0.3, "Virtual Network  10.20.0.0/16", size=11, color=MIDBLUE, bold=True)

# Main ACA env
add_box(s, 4.95, 1.95, 5.2, 2.55, "", fill=WHITE, line=MIDBLUE, line_w=1.25)
add_text(s, 5.1, 2.0, 5.0, 0.3, "Main Container Apps env  (aca-subnet 10.20.2.0/23)", size=10, color=MIDBLUE, bold=True)
add_box(s, 5.12, 2.42, 1.55, 0.62, "learner-portal", fill=LIGHT2, line=AZURE, size=10, color=NAVY)
add_box(s, 6.77, 2.42, 1.55, 0.62, "course-api", fill=LIGHT2, line=AZURE, size=10, color=NAVY)
add_box(s, 8.42, 2.42, 1.6, 0.62, "assessment-api", fill=LIGHT2, line=AZURE, size=10, color=NAVY)
add_box(s, 5.12, 3.12, 1.55, 0.62, "gradebook-api", fill=LIGHT2, line=AZURE, size=10, color=NAVY)
add_box(s, 6.77, 3.12, 3.25, 1.28, "6 Quiz-Service Lanes\nquiz-appgw · quiz-app · quiz-perf\nquiz-query · quiz-pool · quiz-secret",
        fill=CYAN, line=AZURE, size=10, bold=True, color=NAVY)

# NSG-lane env (left) + reporting-worker VM (right) share the bottom row of the VNet
add_box(s, 4.95, 4.62, 3.15, 1.15, "", fill=WHITE, line=MIDBLUE, line_w=1.25)
add_text(s, 5.08, 4.67, 3.0, 0.3, "2nd ACA env  (nsg-lane 10.20.4.0/23)", size=8.5, color=MIDBLUE, bold=True)
add_box(s, 5.07, 5.05, 2.9, 0.58, "quiz-nsg  (isolated nsg lane)", fill=CYAN, line=AZURE, size=9.5, bold=True, color=NAVY)

add_box(s, 8.25, 4.62, 1.9, 1.15, "Reporting VM", fill=AZURE, line=NAVY, size=10, bold=True, color=WHITE,
        sub="vm-subnet 10.20.6.0/24\nB1s + 8 GB data disk\nSyslog → Azure Monitor", sub_size=7, sub_color=LIGHT2)

# arrows AppGw -> envs
add_arrow(s, 4.35, 3.2, 6.77, 3.5, color=AZURE, width=2.0)
add_arrow(s, 4.35, 3.7, 5.3, 5.0, color=AZURE, width=2.0)

# Data tier (right)
add_box(s, 10.75, 2.05, 2.25, 1.15, "Azure Database\nfor PostgreSQL", fill=AZURE, line=NAVY, size=12,
        bold=True, color=WHITE, sub="zava · zava_query\nrole: app_pool", sub_size=8, sub_color=LIGHT2)
add_box(s, 10.75, 3.45, 2.25, 1.1, "Key Vault", fill=AZURE, line=NAVY, size=12, bold=True, color=WHITE,
        sub="db-password\ndb-password-secretlane\ndb-pool-password", sub_size=8, sub_color=LIGHT2)
add_box(s, 10.75, 4.75, 2.25, 0.75, "Container Registry", fill=MIDBLUE, line=NAVY, size=11, bold=True, color=WHITE)
add_arrow(s, 10.1, 3.6, 10.75, 2.7, color=MIDBLUE, width=1.75)   # apps -> pg
add_arrow(s, 10.1, 3.9, 10.75, 3.9, color=MIDBLUE, width=1.75)   # apps -> kv

# Ops tier (bottom)
add_box(s, 0.45, 5.5, 2.0, 1.05, "Azure Monitor\nAlerts", fill=NAVY, line=CYAN, size=12, bold=True,
        color=WHITE, sub="symptom-only", sub_size=8, sub_color=CYAN)
add_box(s, 2.65, 5.5, 1.85, 1.05, "PagerDuty\nIncidents", fill=NAVY, line=CYAN, size=12, bold=True, color=WHITE)
add_box(s, 4.95, 5.95, 5.2, 0.95, "Azure SRE Agent  —  detect · diagnose · mitigate live · open IaC PR",
        fill=AZURE, line=CYAN, size=13, bold=True, color=WHITE)
add_arrow(s, 1.45, 5.5, 1.45, 4.3, color=NAVY, width=1.5, dashed=True)   # apps emit -> monitor (up)
add_arrow(s, 2.45, 6.0, 2.65, 6.0, color=CYAN, width=1.75)
add_arrow(s, 4.5, 6.2, 4.95, 6.35, color=CYAN, width=1.75)
add_arrow(s, 7.5, 5.95, 7.5, 5.65, color=CYAN, width=1.75, head=True, tail=True)  # agent <-> apps


# ════════════════════════════════════════════════════════════════════
# Slide 4 — Network topology
# ════════════════════════════════════════════════════════════════════
s = blank()
header(s, "Network Topology & Isolation", "One VNet, four subnets, four NSGs — lanes share infra but stay isolated where it matters")
outer = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _in(0.6), _in(1.35), _in(12.1), _in(4.0))
outer.fill.solid(); outer.fill.fore_color.rgb = LIGHT1; outer.line.color.rgb = AZURE; outer.line.width = Pt(1.5)
outer.shadow.inherit = False
add_text(s, 0.8, 1.42, 8, 0.3, "vnet-zava   10.20.0.0/16", size=13, color=MIDBLUE, bold=True)

add_box(s, 0.82, 2.0, 2.78, 3.1, "", fill=WHITE, line=MIDBLUE)
add_text(s, 0.95, 2.07, 2.6, 0.5, "appgw-subnet\n10.20.1.0/24", size=10, color=MIDBLUE, bold=True)
add_box(s, 0.92, 2.7, 2.58, 1.0, "Application Gateway", fill=AZURE, line=NAVY, size=11, bold=True, color=WHITE)
add_box(s, 0.92, 3.85, 2.58, 0.7, "NSG: nsg-appgw", fill=LIGHT2, line=AZURE, size=9.5, color=NAVY,
        sub="allow :80/:443 + :8081–:8087", sub_size=7.5)

add_box(s, 3.73, 2.0, 2.78, 3.1, "", fill=WHITE, line=MIDBLUE)
add_text(s, 3.86, 2.07, 2.6, 0.5, "aca-subnet\n10.20.2.0/23", size=10, color=MIDBLUE, bold=True)
add_box(s, 3.83, 2.7, 2.58, 1.0, "Main ACA env\nportal + APIs + 6 lanes", fill=CYAN, line=AZURE, size=10.5,
        bold=True, color=NAVY)
add_box(s, 3.83, 3.85, 2.58, 0.7, "NSG: nsg-aca", fill=LIGHT2, line=AZURE, size=9.5, color=NAVY,
        sub="allow AppGw → apps", sub_size=7.5)

add_box(s, 6.64, 2.0, 2.78, 3.1, "", fill=WHITE, line=RED)
add_text(s, 6.77, 2.07, 2.6, 0.5, "nsg-lane-subnet\n10.20.4.0/23", size=10, color=RED, bold=True)
add_box(s, 6.74, 2.7, 2.58, 1.0, "2nd ACA env\nquiz-nsg (isolated)", fill=CYAN, line=AZURE, size=10.5,
        bold=True, color=NAVY)
add_box(s, 6.74, 3.85, 2.58, 0.7, "NSG: nsg-nsglane", fill=REDFILL, line=RED, size=9.5, color=RED,
        sub="nsg scenario breaks ONLY here", sub_size=7.5, sub_color=RED)

add_box(s, 9.55, 2.0, 2.78, 3.1, "", fill=WHITE, line=MIDBLUE)
add_text(s, 9.68, 2.07, 2.6, 0.5, "vm-subnet\n10.20.6.0/24", size=10, color=MIDBLUE, bold=True)
add_box(s, 9.65, 2.7, 2.58, 1.0, "Reporting-worker VM\nnightly grade exports", fill=AZURE, line=NAVY, size=10.5,
        bold=True, color=WHITE, sub="B1s + 8 GB data disk", sub_size=8, sub_color=LIGHT2)
add_box(s, 9.65, 3.85, 2.58, 0.7, "NSG: nsg-vm", fill=LIGHT2, line=AZURE, size=9.5, color=NAVY,
        sub="closed; reached via az run-command", sub_size=7.5)

add_arrow(s, 3.6, 3.2, 3.73, 3.2, color=AZURE, width=2.0)
add_arrow(s, 3.6, 3.0, 6.64, 3.0, color=AZURE, width=2.0)

add_box(s, 0.6, 5.55, 12.1, 1.3,
        "Why split subnets?  The nsg scenario injects a network DENY rule — putting quiz-nsg in its own subnet + NSG\n"
        "blackholes ONLY its lane while the other six keep serving. The reporting-worker VM sits in its own subnet (no\n"
        "inbound; chaos and the agent reach it via az vm run-command). PostgreSQL is reached over VNet-integrated egress.",
        fill=NAVY, line=None, size=12, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════
# Slide 5 — Incident lifecycle
# ════════════════════════════════════════════════════════════════════
s = blank()
header(s, "The Incident Lifecycle", "Every scenario follows the same closed loop — a real break and a durable fix")
steps = [
    ("1  Break", "A real fault ships:\nIaC param, image,\nDB state, or secret.", AZURE),
    ("2  Symptom", "Students feel it.\nSymptom-only —\nno cause named.", MIDBLUE),
    ("3  Detect", "Azure Monitor alert\n+ synthetic monitor\npage PagerDuty.", AZURE),
    ("4  Engage", "SRE Agent acks the\nincident & begins\ndiagnosis.", MIDBLUE),
    ("5  Mitigate", "Agent restores\nservice LIVE\n(fast recovery).", AZURE),
    ("6  Durable Fix", "Agent opens an\nIaC PR to fix the\nroot cause.", MIDBLUE),
]
x = 0.55
for i, (t, d, c) in enumerate(steps):
    add_box(s, x, 2.2, 1.85, 1.9, t, fill=c, line=NAVY, size=15, bold=True, color=WHITE,
            sub="\n" + d, sub_size=10, sub_color=LIGHT2)
    if i < len(steps) - 1:
        add_arrow(s, x + 1.85, 3.15, x + 2.07, 3.15, color=NAVY, width=2.0)
    x += 2.07
add_arrow(s, 12.0, 4.1, 1.5, 4.1, color=GREEN, width=2.0, dashed=True)
add_text(s, 4.0, 4.15, 5.3, 0.4, "reset.ps1 restores baseline for the next run", size=11, color=GREEN, bold=True,
         align=PP_ALIGN.CENTER)
add_box(s, 0.55, 5.0, 12.2, 1.7,
        "Two detection paths:  log-based Azure Monitor alerts catch the latency / error scenarios, while a synthetic\n"
        "student-journey monitor pages PagerDuty directly for connectivity blackholes that emit no logs.  The simulator\n"
        "(simulator/demo.py) narrates the whole loop and deep-links to the live SRE Agent investigation thread.",
        fill=LIGHT1, line=AZURE, size=12.5, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════
# Slide 6 — The 7 lanes map
# ════════════════════════════════════════════════════════════════════
s = blank()
header(s, "Seven Parallel Fault Lanes", "Same App Gateway, one frontend port per lane — all scenarios run at once")
lanes = [
    ("8081", "nsg", "quiz-nsg", "Network DENY rule", "2nd ACA env"),
    ("8082", "appgw", "quiz-appgw", "Bad health-probe path", "App Gateway"),
    ("8083", "app", "quiz-app", "Scaled to zero replicas", "Container Apps"),
    ("8084", "perf", "quiz-perf", "Slow code release (v1.1 image)", "App code"),
    ("8085", "query", "quiz-query", "Question-bank index corruption", "PostgreSQL zava_query"),
    ("8086", "pool", "quiz-pool", "DB role connection limit", "PostgreSQL role"),
    ("8087", "secret", "quiz-secret", "Rotated invalid credential", "Key Vault"),
]
add_box(s, 0.55, 1.5, 2.2, 4.8, "Application\nGateway\n\npublic IP", fill=AZURE, line=NAVY, size=15,
        bold=True, color=WHITE, sub=":80 → learner-portal", sub_size=9, sub_color=LIGHT2)
y = 1.5
for port, key, app, fault, tier in lanes:
    add_box(s, 3.4, y, 1.05, 0.62, ":" + port, fill=NAVY, line=CYAN, size=13, bold=True, color=WHITE)
    add_box(s, 4.7, y, 2.4, 0.62, app, fill=LIGHT2, line=AZURE, size=12, bold=True, color=NAVY)
    add_box(s, 7.35, y, 3.4, 0.62, fault, fill=REDFILL, line=RED, size=11, color=RED)
    add_box(s, 11.0, y, 1.75, 0.62, tier, fill=LIGHT1, line=MIDBLUE, size=9, color=MIDBLUE)
    add_arrow(s, 2.75, y + 0.31, 3.4, y + 0.31, color=AZURE, width=1.5)
    add_arrow(s, 4.45, y + 0.31, 4.7, y + 0.31, color=AZURE, width=1.5)
    y += 0.69
add_text(s, 3.4, 6.45, 9, 0.4, "lane (operator)        container app           planted fault (symptom-only to students)        breaks at",
         size=9, color=GRAY, italic=True)


# ════════════════════════════════════════════════════════════════════
# Slide 7 — How the SRE Agent is wired (shared across every scenario)
# ════════════════════════════════════════════════════════════════════
s = blank()
header(s, "The SRE Agent — How It's Wired", "The incident-responder handles every Zava incident from telemetry; 3 scheduled auditors run weekly governance (next section)")
add_box(s, 0.55, 1.45, 5.95, 1.2,
        "Incident agent:  zava-incident-responder",
        fill=NAVY, line=None, size=14, bold=True, color=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
add_text(s, 0.72, 1.95, 5.6, 0.7,
         "ExtendedAgent · autonomous · temperature 0.2 · critic-on-handoff gate.\n"
         "Applied idempotently via scripts/configure-agent (never by hand).",
         size=10.5, color=WHITE)
add_box(s, 6.8, 1.45, 5.95, 1.2,
        "Incident filter:  zava-learning-response",
        fill=NAVY, line=None, size=14, bold=True, color=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
add_text(s, 6.97, 1.95, 5.6, 0.7,
         "title contains 'Zava' · Sev0–Sev4 · agentMode = autonomous ·\n"
         "handlingAgent → the agent · ≤3 attempts · 3 h merge window.",
         size=10.5, color=WHITE)
add_box(s, 0.55, 2.95, 12.2, 0.42, "Skills it can load (Zava + built-in system skills)",
        fill=AZURE, line=None, size=12.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, 0.7, 3.5, 12.0, 1.5,
         "Triage:   connectivity-triage  (edge / network)      ·      performance-investigation  (app tier + reporting VM)\n"
         "Then, in order:   rca-analysis → evidence-before-after → recommendations-next-steps → pr-delivery →\n"
         "                          servicenow-change-management → zava-reporting → pagerduty-incident-update\n"
         "System skills (diagnostic support):   network-connectivity · load-balancer · application-gateway · network-topology",
         size=11.5, color=INK)
panel(s, 0.55, 5.15, 6.0, 1.75, "Connectors  (data it reads)", [
    "•  Application Insights",
    "•  Log Analytics",
    "•  Azure Monitor",
    "•  Microsoft Learn  (MCP)",
], body_size=11)
panel(s, 6.75, 5.15, 6.0, 1.75, "Tools  (actions it takes)", [
    "•  RunAzCli read + write commands",
    "•  Python — builds the HTML report & audit decks on disk",
    "•  CreateServiceNowChangeRequest · UploadServiceNowAttachment",
    "•  SearchMemory · log_analytics_query",
], body_size=11)


# ════════════════════════════════════════════════════════════════════
# Slide 8 — What each skill does (shared catalog)
# ════════════════════════════════════════════════════════════════════
s = blank()
header(s, "What Each Skill Does",
       "Incident-response skills: the agent picks ONE triage skill from telemetry, then runs the rest in order")
cat = [
    ("connectivity-triage", MIDBLUE, [
        "TRIAGE · edge / network",
        "Traces App Gateway → NSG →",
        "Container Apps → APIs, finds the",
        "broken hop, and fixes it.",
        "Used by: nsg · appgw · secret",
    ]),
    ("performance-investigation", MIDBLUE, [
        "TRIAGE · app / worker tier",
        "Reachable but degraded — 5xx,",
        "slow, zero instances, or a VM",
        "batch job failing. Fixes the tier.",
        "Used by: app·perf·query·pool·disk",
    ]),
    ("pagerduty-incident-update", AZURE, [
        "Acknowledges the incident, posts",
        "symptom-only status notes, and",
        "resolves it once recovery is",
        "verified.",
    ]),
    ("rca-analysis", AZURE, [
        "Root-cause narrative: incident",
        "timeline, trigger vs. latent cause",
        "(5-Whys), and contributing",
        "factors.",
    ]),
    ("evidence-before-after", AZURE, [
        "Renders the ONE before/after",
        "visual (path/topology or time-",
        "series) plus a delta table that",
        "proves impact and recovery.",
    ]),
    ("recommendations-next-steps", AZURE, [
        "Prioritized preventive, detective,",
        "and process actions — with owners,",
        "target dates, and the risk of",
        "inaction.",
    ]),
    ("pr-delivery", AZURE, [
        "Opens the GitHub pull request for",
        "the durable fix — Bicep (infra/) or",
        "application code (src/). The PR the",
        "Change Request references.",
    ]),
    ("servicenow-change-management", AZURE, [
        "Raises a ServiceNow Change",
        "Request referencing the PR and",
        "attaches the RCA report.",
        "(Best-effort — never blocks.)",
    ]),
    ("zava-reporting", AZURE, [
        "Branded in-thread executive",
        "summary + a downloadable HTML",
        "report; assembles the other",
        "skills' output.",
    ]),
]
cx = [0.45, 4.45, 8.45]
cw, ch = 3.85, 1.62
cy0 = 1.55
for i, (nm, accent, lines) in enumerate(cat):
    col = i % 3
    row = i // 3
    x = cx[col]
    y = cy0 + row * 1.72
    panel(s, x, y, cw, ch, nm, lines, title_fill=accent, body_size=9)
add_text(s, 0.45, 6.95, 12.4, 0.5,
         "Built-in system skills (network-connectivity · load-balancer · application-gateway · network-topology) reinforce triage. "
         "Audit skills (nsg-audit · rbac-audit · cost-analysis · zava-audit-report) + the always-on redaction-guard → see Proactive Governance Audits.",
         size=9.5, color=GRAY, italic=True)


# ════════════════════════════════════════════════════════════════════
# Slide 9 — Response runbook → artifacts (shared across every scenario)
# ════════════════════════════════════════════════════════════════════
s = blank()
header(s, "Response Runbook → Artifacts", "Every incident follows the same ordered runbook and must produce every artifact")
runbook = [
    ("0  Acknowledge", "PagerDuty ack so\nescalation stops."),
    ("1  Triage", "Edge/network vs.\napp tier from\ntelemetry."),
    ("2  Mitigate", "Smallest live fix\n(recorded as a\nnote)."),
    ("3  RCA", "Evidence-backed\nroot-cause\nnarrative."),
    ("4  Evidence", "Before/after data\nproving impact +\nrecovery."),
    ("5  Recommend", "Prioritized\npreventive / next\nsteps."),
    ("6  PR", "GitHub IaC/code\npull request\n(durable fix)."),
    ("7  Change Mgmt", "ServiceNow CR\nref PR (best-\neffort)."),
    ("8  Report", "Branded\ndownloadable\nHTML report."),
    ("9  Resolve", "PD summary note\n+ resolve once\nverified."),
]
x = 0.4
for i, (t, d) in enumerate(runbook):
    c = AZURE if i % 2 == 0 else MIDBLUE
    add_box(s, x, 1.55, 1.2, 1.85, t, fill=c, line=NAVY, size=11, bold=True, color=WHITE,
            sub="\n" + d, sub_size=8.5, sub_color=LIGHT2, anchor=MSO_ANCHOR.TOP)
    if i < len(runbook) - 1:
        add_arrow(s, x + 1.2, 2.45, x + 1.29, 2.45, color=NAVY, width=1.5)
    x += 1.29
add_box(s, 0.55, 3.95, 12.2, 0.42, "Artifacts the agent must produce (critic-on-handoff blocks completion until all exist)",
        fill=AZURE, line=None, size=12.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
panel(s, 0.55, 4.55, 6.0, 2.25, "Investigation & evidence", [
    "•  PagerDuty incident — acknowledged then resolved",
    "•  Root-cause analysis narrative",
    "•  Before/after evidence (data proving recovery)",
    "•  Prioritized recommendations & next steps",
], body_size=11)
panel(s, 6.75, 4.55, 6.0, 2.25, "Durable fix & hand-off", [
    "•  GitHub pull request (IaC or code) for the durable fix",
    "•  ServiceNow Change Request referencing the PR + RCA",
    "•  Downloadable branded HTML executive report",
    "•  Symptom-only PD closure note (cause never assumed)",
], body_size=11)


# ════════════════════════════════════════════════════════════════════
# Per-scenario good vs bad slides
# ════════════════════════════════════════════════════════════════════
def chain_slide(title, sub, port, nodes, broke_idx, good_caption, bad_caption,
                symptom, root_cause, mechanism):
    s = blank()
    header(s, title, sub)
    # symptom / root-cause / mechanism strip
    add_box(s, 0.55, 1.55, 4.05, 0.9, "Student symptom", fill=NAVY, line=None, size=11, bold=True,
            color=CYAN, sub=symptom, sub_size=10, sub_color=WHITE, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT)
    add_box(s, 4.75, 1.55, 4.05, 0.9, "Root cause (agent diagnoses)", fill=NAVY, line=None, size=11,
            bold=True, color=CYAN, sub=root_cause, sub_size=10, sub_color=WHITE, anchor=MSO_ANCHOR.TOP,
            align=PP_ALIGN.LEFT)
    add_box(s, 8.95, 1.55, 3.8, 0.9, "Fault mechanism", fill=NAVY, line=None, size=11, bold=True,
            color=CYAN, sub=mechanism, sub_size=10, sub_color=WHITE, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT)

    def render(x0, w, state, broken):
        col = GREEN if state == "GOOD" else RED
        fillbar = GREENFILL if state == "GOOD" else REDFILL
        add_box(s, x0, 2.65, w, 0.5, state + " STATE", fill=col, line=None, size=14, bold=True, color=WHITE)
        ny = 3.25
        step = 0.86
        nh = 0.6
        bw = w - 1.0
        bx = x0 + 0.5
        for i, (label, sub_t) in enumerate(nodes):
            is_broken = broken and i == broke_idx
            nfill = REDFILL if is_broken else (LIGHT2 if i else AZURE)
            nline = RED if is_broken else AZURE
            ncolor = RED if is_broken else (NAVY if i else WHITE)
            add_box(s, bx, ny, bw, nh, label, fill=nfill, line=nline, size=11.5, bold=(i == 0),
                    color=ncolor, sub=sub_t, sub_size=8.5,
                    sub_color=(RED if is_broken else (LIGHT2 if i == 0 else GRAY)))
            if i < len(nodes) - 1:
                edge_broken = broken and (i + 1 == broke_idx)
                acol = RED if edge_broken else col
                add_arrow(s, bx + bw / 2, ny + nh, bx + bw / 2, ny + step, color=acol, width=2.0)
                if edge_broken:
                    add_text(s, bx + bw / 2 + 0.05, ny + nh - 0.02, 0.6, 0.3, "X", size=16, color=RED, bold=True)
            ny += step
        cap = good_caption if state == "GOOD" else bad_caption
        add_box(s, x0, 6.6, w, 0.7, cap, fill=fillbar, line=col, size=10.5,
                color=(GREEN if state == "GOOD" else RED), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    render(0.55, 6.0, "GOOD", False)
    render(6.95, 6.0, "BAD", True)
    return s


# Per-scenario response detail (Approach 1): a light pictorial pipeline.
# Fault -> Trigger -> Agent -> the ONE triage skill specific to this scenario.
# The shared close-out skills are documented once on the "What Each Skill Does" slide.
TRIAGE_WHAT = {
    "connectivity-triage":
        "Traces the request path App Gateway \u2192 NSG \u2192 Container Apps \u2192 APIs,\n"
        "finds the broken hop, and remediates it.",
    "performance-investigation":
        "Platform is reachable but a tier is degraded \u2014 diagnoses from App\n"
        "Insights / Log Analytics / Syslog and fixes the app or worker tier.",
}

RESPONSE = {
    "nsg": {
        "title": "Scenario 1 \u2014 Response (NSG)",
        "triage": "connectivity-triage",
        "fault": "NSG legacy DENY rule shipped\nfrom IaC (priority-100 DENY\nshadows the ALLOW).",
        "trigger": "Synthetic student monitor \u2192\nPagerDuty:\n'quiz launches failing'.",
        "why": "Edge/network symptom \u2014 the\ngateway cannot reach quiz-nsg.",
        "mitig": "remove the legacy DENY rule",
        "durable": "revert injectLegacyDenyNsgLane (infra/)",
    },
    "appgw": {
        "title": "Scenario 2 \u2014 Response (App Gateway)",
        "triage": "connectivity-triage",
        "fault": "App Gateway probe path flipped\nto /status-ping from IaC \u2192\nbackend marked unhealthy (502s).",
        "trigger": "Azure Monitor 5xx alert \u2192\nPagerDuty:\n'portal returning 502s'.",
        "why": "Edge symptom \u2014 probe hits a\nmissing path, no healthy backend.",
        "mitig": "repoint the probe to /health",
        "durable": "restore appgwLaneProbePath=/health (infra/)",
    },
    "app": {
        "title": "Scenario 3 \u2014 Response (App)",
        "triage": "performance-investigation",
        "fault": "Replica bounds set to 0 from\nIaC \u2192 quiz-app scaled to 0/0\n(no instances).",
        "trigger": "Monitor \u2192 PagerDuty:\n'quiz service unavailable'.",
        "why": "Platform reachable; the app tier\nhas zero healthy instances.",
        "mitig": "scale quiz-app back up",
        "durable": "restore replica bounds 1\u20133 (infra/)",
    },
    "perf": {
        "title": "Scenario 4 \u2014 Response (Perf)",
        "triage": "performance-investigation",
        "fault": "Real v1.1 image with a\nsynchronous pbkdf2 KDF built\n& rolled out to quiz-perf.",
        "trigger": "Azure Monitor latency alert \u2192\nPagerDuty:\n'quiz launches slow'.",
        "why": "App tier degraded \u2014 heavy crypto\non the hot quiz path.",
        "mitig": "roll back to the v1.0 revision",
        "durable": "roll back the v1.1 code release (src/)",
    },
    "query": {
        "title": "Scenario 5 \u2014 Response (Query)",
        "triage": "performance-investigation",
        "fault": "question_bank index dropped on\nzava_query (simulated\ncorruption) \u2192 500k-row scan.",
        "trigger": "Azure Monitor latency alert \u2192\nPagerDuty:\n'quiz loading slowly'.",
        "why": "App reachable; the database does\na full table scan every request.",
        "mitig": "REINDEX / rebuild the index",
        "durable": "codify the index in schema/migration",
    },
    "pool": {
        "title": "Scenario 6 \u2014 Response (Pool)",
        "triage": "performance-investigation",
        "fault": "app_pool DB role capped to 1\nconnection (live DB drift).",
        "trigger": "Azure Monitor error-rate alert \u2192\nPagerDuty:\n'errors under load'.",
        "why": "Reachable, but the pool starves\nunder concurrency \u2192 500s.",
        "mitig": "reset the role connection limit",
        "durable": "codify the role connection policy",
    },
    "secret": {
        "title": "Scenario 7 \u2014 Response (Secret)",
        "triage": "connectivity-triage",
        "fault": "Key Vault secret db-password-\nsecretlane rotated to an invalid\nvalue (live drift).",
        "trigger": "Azure Monitor auth-failure alert \u2192\nPagerDuty:\n'service failing'.",
        "why": "The service can no longer\nauthenticate to its database.",
        "mitig": "restore the valid Key Vault secret",
        "durable": "add secret-rotation guardrails",
    },
    "disk": {
        "title": "Scenario 8 \u2014 Response (Reporting VM)",
        "triage": "performance-investigation",
        "fault": "Reporting-worker VM /data filled\nvia az vm run-command (live VM\ndrift) \u2192 exports can't write.",
        "trigger": "Azure Monitor Syslog alert:\nZava-grade-exports-failing\n(no synthetic monitor).",
        "why": "Back-office VM batch job fails;\nthe student site stays healthy.",
        "mitig": "free space on /data",
        "durable": "add disk-pressure guardrails to IaC",
    },
}

RESPONSE_SUB = "Fault \u2192 trigger \u2192 agent \u2192 the one triage skill specific to this scenario"


def response_slide(key):
    d = RESPONSE[key]
    s = blank()
    header(s, d["title"], RESPONSE_SUB)
    bx = [0.5, 3.65, 6.8, 9.95]
    bw, by, bh = 2.85, 1.75, 2.7
    add_box(s, bx[0], by, bw, bh, "\u2460  Fault", fill=REDFILL, line=RED, size=14, bold=True,
            color=RED, anchor=MSO_ANCHOR.TOP, sub="\n" + d["fault"], sub_size=10.5, sub_color=INK)
    add_box(s, bx[1], by, bw, bh, "\u2461  Trigger", fill=LIGHT1, line=AMBER, size=14, bold=True,
            color=AMBER, anchor=MSO_ANCHOR.TOP, sub="\n" + d["trigger"], sub_size=10.5, sub_color=INK)
    add_box(s, bx[2], by, bw, bh, "\u2462  Agent", fill=NAVY, line=None, size=14, bold=True,
            color=CYAN, anchor=MSO_ANCHOR.TOP,
            sub="\nzava-incident-responder\n\nautonomous \u00b7 \u22643 attempts\n(routed by the incident filter)",
            sub_size=10.5, sub_color=WHITE)
    add_box(s, bx[3], by, bw, bh, "\u2463  " + d["triage"], fill=LIGHT2, line=AZURE, line_w=2.5,
            size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.TOP,
            sub="\n" + TRIAGE_WHAT[d["triage"]] + "\n\nWhy here:  " + d["why"],
            sub_size=9.5, sub_color=INK)
    for i in range(3):
        add_arrow(s, bx[i] + bw, by + bh / 2, bx[i + 1], by + bh / 2, color=NAVY, width=2.0)
    add_text(s, bx[3], by - 0.34, bw, 0.3, "\u2605 skill specific to this scenario", size=9.5,
             color=AZURE, bold=True, align=PP_ALIGN.CENTER)
    add_box(s, 0.5, 4.75, 12.3, 1.0,
            "Then the agent runs the shared close-out runbook (same for every scenario)",
            fill=AZURE, line=None, size=12.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP,
            sub="rca-analysis \u2192 evidence-before-after \u2192 recommendations-next-steps \u2192 pr-delivery \u2192 "
                "servicenow-change-management \u2192 zava-reporting \u2192 pagerduty-incident-update"
                "          \u25b8 see \u201cWhat Each Skill Does\u201d",
            sub_size=10.5, sub_color=WHITE)
    add_box(s, 0.5, 5.9, 12.3, 0.78, "Artifacts", fill=LIGHT1, line=AZURE, size=11.5, bold=True,
            color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            sub="live mitigation (" + d["mitig"] + ")  \u00b7  GitHub PR (" + d["durable"] +
                ")  \u00b7  ServiceNow CR  \u00b7  RCA  \u00b7  before/after evidence  \u00b7  HTML report  \u00b7  PagerDuty resolved",
            sub_size=10, sub_color=INK)
    add_text(s, 0.5, 6.82, 12.3, 0.4,
             "Connectors the agent reads:  App Insights \u00b7 Log Analytics \u00b7 Azure Monitor \u00b7 Microsoft Learn (MCP)",
             size=10, color=GRAY, italic=True)
    return s


# Common entry nodes
def nset(*extra):
    return [("Student", "browser"), (f"Application Gateway", "")] + list(extra)


# ════════════════════════════════════════════════════════════════════
# Proactive Governance Audits — section (after the shared incident
# "how the agent works" block, before the per-scenario deep dives)
# ════════════════════════════════════════════════════════════════════

# ── Audit slide A — two ways the agent works ────────────────────────
s = blank()
header(s, "Two Ways the SRE Agent Works",
       "The same agent platform runs reactive incident response AND proactive weekly governance audits")
add_box(s, 0.55, 1.45, 6.0, 0.55, "REACTIVE  —  incident response", fill=AZURE, line=None,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_box(s, 0.55, 2.05, 6.0, 4.45, "", fill=LIGHT1, line=AZURE, line_w=1.5)
add_text(s, 0.75, 2.18, 5.65, 4.2,
         "Trigger:  an Azure Monitor / synthetic alert pages\n"
         "    PagerDuty.\n"
         "Routing:  incident filter zava-learning-response →\n"
         "    zava-incident-responder  (autonomous).\n"
         "Action:  acknowledge → triage → MITIGATE LIVE →\n"
         "    RCA → evidence → recommendations.\n"
         "Hand-off:  GitHub IaC / code PR (durable fix) ·\n"
         "    ServiceNow CR · branded HTML report ·\n"
         "    PagerDuty resolved.\n"
         "Cadence:  on demand, per incident.\n"
         "Mode:  WRITES live mitigations + opens the fix PR.",
         size=11.5, color=INK)
add_box(s, 6.8, 1.45, 6.0, 0.55, "PROACTIVE  —  scheduled governance", fill=GREEN, line=None,
        size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_box(s, 6.8, 2.05, 6.0, 4.45, "", fill=GREENFILL, line=GREEN, line_w=1.5)
add_text(s, 7.0, 2.18, 5.65, 4.2,
         "Trigger:  a weekly cron Scheduled Task — no\n"
         "    incident, no PagerDuty.\n"
         "Routing:  each task → its own auditor agent\n"
         "    (zava-nsg / rbac / cost).\n"
         "Action:  enumerate & assess posture, score every\n"
         "    finding SEV1 / SEV2 / SEV3.\n"
         "Hand-off:  ONE branded, downloadable PowerPoint\n"
         "    deck + a short posture summary with a\n"
         "    clickable download link.\n"
         "Cadence:  weekly, per audit type.\n"
         "Mode:  READ-ONLY — recommend, never change.",
         size=11.5, color=INK)
add_box(s, 0.55, 6.68, 12.25, 0.62,
        "Same custom-agent platform, connectors, and always-on redaction guard — all applied idempotently via "
        "scripts/configure-agent (never by hand).",
        fill=NAVY, line=None, size=11.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# ── Audit slide B — the three weekly audits ─────────────────────────
s = blank()
header(s, "Proactive Governance Audits",
       "Three read-only weekly audits — each its own Scheduled Task, agent, and skill; each delivers a branded deck")
acols = [
    (0.55, 1.85, "Schedule"),
    (2.45, 2.75, "Scheduled task"),
    (5.25, 2.25, "Auditor agent"),
    (7.55, 1.85, "Skill"),
    (9.45, 3.35, "What it inspects"),
]
for ax, aw, label in acols:
    add_box(s, ax, 1.5, aw, 0.5, label, fill=NAVY, line=None, size=11, bold=True, color=CYAN)
arows = [
    ("Mon 08:00\n0 8 * * 1", "zava-nsg-weekly-audit", "zava-nsg-auditor", "nsg-audit",
     "NSGs & rules, route tables / UDRs, Private Endpoints & DNS, publicly-exposed data — vs. the "
     "intended :8081–:8087 lanes"),
    ("Tue 08:00\n0 8 * * 2", "zava-rbac-weekly-audit", "zava-rbac-auditor", "rbac-audit",
     "Role assignments at RG scope + Activity-Log usage — over-privileged, direct, guest, stale, and "
     "privileged-but-unused access"),
    ("Wed 08:00\n0 8 * * 3", "zava-cost-weekly-analysis", "zava-cost-analyst", "cost-analysis",
     "Cost Management WoW spend per resource — top drivers, Δ% trend, and idle / oversized / orphaned "
     "resources with savings"),
]
ry = 2.05
rh = 1.28
for sched, task, agent, skill, inspects in arows:
    cells = [
        (acols[0][0], acols[0][1], sched, LIGHT2, NAVY, 10.5, True),
        (acols[1][0], acols[1][1], task, WHITE, MIDBLUE, 10.5, True),
        (acols[2][0], acols[2][1], agent, WHITE, NAVY, 10.5, False),
        (acols[3][0], acols[3][1], skill, WHITE, AZURE, 10.5, False),
        (acols[4][0], acols[4][1], inspects, LIGHT1, INK, 9.5, False),
    ]
    for cx, cw, txt, fill, col, sz, bold in cells:
        align = PP_ALIGN.CENTER if cw < 3.0 else PP_ALIGN.LEFT
        add_box(s, cx, ry, cw, rh, txt, fill=fill, line=AZURE, size=sz, bold=bold, color=col,
                align=align, anchor=MSO_ANCHOR.MIDDLE)
    ry += rh + 0.08
add_box(s, 0.55, 6.05, 12.25, 1.2,
        "Every audit hands its findings to the zava-audit-report skill",
        fill=AZURE, line=None, size=12.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
        sub="→ ONE branded PowerPoint (posture summary · SEV1/2/3 findings table · recommendations) "
            "verified on disk, with a clickable download link.   READ-ONLY: every auditor recommends, never changes.",
        sub_size=10.5, sub_color=WHITE)

# ── Audit slide C — how a weekly audit runs ─────────────────────────
s = blank()
header(s, "How a Weekly Audit Runs",
       "Every auditor follows the same 3-step read-only runbook and must produce the branded deck")
astep = [
    ("1  Audit / Analyse",
     "Load the audit skill (nsg-audit /\nrbac-audit / cost-analysis):\nenumerate & assess posture,\n"
     "score each finding SEV1/2/3,\nbuild findings + recommendations.", AZURE),
    ("2  Report",
     "Hand the findings to the\nzava-audit-report skill — it\nrenders ONE branded PowerPoint\n"
     "(posture summary, findings table,\ncost-driver chart) into the thread.", MIDBLUE),
    ("3  Deliver",
     "mkdir → write → VERIFY on disk\n(non-zero bytes) → post a short\nposture summary + the deck's\n"
     "/api/files/… link as a clickable\nmarkdown hyperlink (never a table).", AZURE),
]
ax = 0.55
aw, ay, ah = 3.9, 1.55, 2.55
for i, (t, d, c) in enumerate(astep):
    add_box(s, ax, ay, aw, ah, t, fill=c, line=NAVY, size=15, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.TOP, sub="\n" + d, sub_size=10.5, sub_color=LIGHT2)
    if i < len(astep) - 1:
        add_arrow(s, ax + aw, ay + ah / 2, ax + aw + 0.3, ay + ah / 2, color=NAVY, width=2.0)
    ax += aw + 0.3
panel(s, 0.55, 4.55, 3.9, 2.4, "Read-only by design", [
    "Auditors enumerate and recommend —",
    "they NEVER create, modify, or delete",
    "an NSG rule, role assignment, or",
    "resource. No incident, no PagerDuty.",
], title_fill=GREEN, title_color=WHITE, body_size=10)
panel(s, 4.75, 4.55, 3.9, 2.4, "Redaction always on", [
    "Before any thread message or deck",
    "leaves the agent, the redaction-guard",
    "scrubber (SearchMemory \"zava-",
    "redaction\") strips secrets / tokens /",
    "PII. Names, rules, roles, SKUs, $ kept.",
], title_fill=NAVY, body_size=10)
panel(s, 8.95, 4.55, 3.85, 2.4, "Severity model", [
    "SEV1 = act now · SEV2 = this week ·",
    "SEV3 = tidy-up. Intended lab exposure",
    "(AppGw :8081–:8087) and first-party",
    "platform identities are recognised and",
    "listed, not flagged as findings.",
], title_fill=AZURE, body_size=10)


chain_slide(
    "Scenario 1 — Connectivity (NSG)", "Lane :8081 · quiz-nsg in the isolated 2nd Container Apps env",
    "8081",
    [("Student", "launches a quiz"),
     ("Application Gateway  :8081", "routes to the nsg lane"),
     ("NSG  nsg-nsglane", "AppGw → apps allowed"),
     ("quiz-nsg", "serves quiz content")],
    broke_idx=2,
    good_caption="NSG ships clean: the ALLOW rule (priority 200) lets the gateway reach quiz-nsg.",
    bad_caption="A legacy 'segmentation' DENY at priority 100 beats the ALLOW — traffic is blackholed.",
    symptom="“Launch quiz” spins and fails; nothing was deployed today.",
    root_cause="NSG priority inversion: a high-priority DENY shadows the ALLOW rule.",
    mechanism="IaC param injectLegacyDenyNsgLane = true\n(committed + live NSG rule).")
response_slide("nsg")

chain_slide(
    "Scenario 2 — Gateway Backend (App Gateway)", "Lane :8082 · quiz-appgw health probe",
    "8082",
    [("Student", "loads the portal"),
     ("Application Gateway  :8082", "health probe → backend"),
     ("Health probe path", "GET /health → 200"),
     ("quiz-appgw", "healthy backend")],
    broke_idx=2,
    good_caption="Probe path /health returns 200, so the gateway keeps the backend in rotation.",
    bad_caption="Probe path points at a missing page → every backend marked unhealthy → 502s.",
    symptom="The portal returns 502 errors on every page; apps look up.",
    root_cause="App Gateway health-probe path misconfigured to a path the app doesn't serve.",
    mechanism="IaC param appgwLaneProbePath\n= /status-ping (committed + live probe).")
response_slide("appgw")

chain_slide(
    "Scenario 3 — No Healthy Instances (App)", "Lane :8083 · quiz-app replica bounds",
    "8083",
    [("Student", "launches a quiz"),
     ("Application Gateway  :8083", "routes to the app lane"),
     ("quiz-app  replicas", "min 1 / max 3"),
     ("Quiz content", "served")],
    broke_idx=2,
    good_caption="Replica bounds 1–3 keep at least one instance answering requests.",
    bad_caption="Replicas pinned to 0 — nothing is left to serve quiz content; launches fail.",
    symptom="Quiz launches fail though the network looks clean.",
    root_cause="The assessment service has zero healthy instances (scaled to zero).",
    mechanism="IaC params appLaneMin/MaxReplicas = 0\n(committed + live scale).")
response_slide("app")

chain_slide(
    "Scenario 4 — Slow Code Release (Perf)", "Lane :8084 · quiz-perf container image",
    "8084",
    [("Student", "launches a quiz"),
     ("Application Gateway  :8084", "routes to the perf lane"),
     ("quiz-perf image", "v1.0 clean build"),
     ("PostgreSQL  zava", "fast response")],
    broke_idx=2,
    good_caption="The clean v1.0 image returns quiz content in a few milliseconds.",
    bad_caption="v1.1 adds a synchronous KDF (pbkdf2, 1.2M iters) on the hot path → 1–2 s per request.",
    symptom="Quizzes load but crawl; latency climbed right after the morning release.",
    root_cause="A real bad code release runs heavy synchronous crypto on the request path.",
    mechanism="Real source → real image: releases/server.v1.1.js\nbuilt & rolled out to quiz-perf.")
response_slide("perf")

chain_slide(
    "Scenario 5 — Database Query (Query)", "Lane :8085 · quiz-query on isolated DB zava_query",
    "8085",
    [("Student", "launches a quiz"),
     ("Application Gateway  :8085", "routes to the query lane"),
     ("quiz-query", "counts the question bank"),
     ("PostgreSQL  zava_query", "index healthy")],
    broke_idx=3,
    good_caption="idx_question_bank_course makes the 500k-row lookup an index scan (fast).",
    bad_caption="Index corrupt/unusable → full sequential scan of 500k rows on every quiz → seconds.",
    symptom="Opening a quiz takes several seconds; app + release look clean.",
    root_cause="A corrupt, unusable index forces a full table scan of the question bank.",
    mechanism="question_bank index corruption on\nzava_query; fix = REINDEX/rebuild.")
response_slide("query")

chain_slide(
    "Scenario 6 — Connection Exhaustion (Pool)", "Lane :8086 · quiz-pool via DB role app_pool",
    "8086",
    [("Student (under load)", "launches a quiz"),
     ("Application Gateway  :8086", "routes to the pool lane"),
     ("quiz-pool  (role app_pool)", "opens DB connections"),
     ("PostgreSQL  zava", "CONNECTION LIMIT -1")],
    broke_idx=3,
    good_caption="Role app_pool has no connection cap, so the pool serves concurrent students.",
    bad_caption="ALTER ROLE app_pool CONNECTION LIMIT 1 → pool starves under load → 500s.",
    symptom="Under exam load some students get errors while others succeed.",
    root_cause="The database role caps connections, starving the service's pool under concurrency.",
    mechanism="Live DB role drift:\nALTER ROLE app_pool CONNECTION LIMIT 1.")
response_slide("pool")

chain_slide(
    "Scenario 7 — Bad Credential (Secret)", "Lane :8087 · quiz-secret via Key Vault secret",
    "8087",
    [("Student", "launches a quiz"),
     ("Application Gateway  :8087", "routes to the secret lane"),
     ("Key Vault secret", "db-password-secretlane"),
     ("PostgreSQL  zava", "auth succeeds")],
    broke_idx=2,
    good_caption="The lane's Key Vault secret holds the real password, so DB auth succeeds.",
    bad_caption="Secret rotated to an invalid value → service can't authenticate → all launches fail.",
    symptom="Every quiz launch on this service fails; nothing was deployed.",
    root_cause="The database credential in Key Vault was rotated to a value that no longer works.",
    mechanism="Live secret drift: rotate Key Vault\ndb-password-secretlane to an invalid value.")
response_slide("secret")

chain_slide(
    "Scenario 8 — Disk Pressure (Reporting Worker)", "Back-office VM · vm-zava-reporting (no AppGw lane)",
    "n/a",
    [("Nightly grade-export job", "writes export files"),
     ("Reporting-worker VM", "runs the export timer"),
     ("Data disk  /data", "stores exports"),
     ("Instructor grade export", "delivered")],
    broke_idx=2,
    good_caption="Exports write to /data and rotate; the job logs a success heartbeat to Syslog.",
    bad_caption="Data disk fills (No space left on device) → every export run fails → no exports.",
    symptom="Instructors stop receiving nightly grade exports; the student site is healthy.",
    root_cause="The reporting worker's data disk filled up, so export writes fail.",
    mechanism="Live VM drift: fill /data on the\nreporting worker; fix = free space.")
response_slide("disk")


# ════════════════════════════════════════════════════════════════════
# Final — how it all connects / legend
# ════════════════════════════════════════════════════════════════════
s = blank()
header(s, "How It All Connects", "One platform, one environment, eight independent stories")
add_box(s, 0.55, 1.4, 12.2, 1.5,
        "Students → Application Gateway → Container Apps → PostgreSQL + Key Vault.  Azure Monitor and a synthetic\n"
        "monitor page PagerDuty; the SRE Agent acknowledges, diagnoses against this exact topology, mitigates live,\n"
        "and opens an IaC pull request for the durable fix. reset.ps1 returns every lane to baseline.",
        fill=LIGHT1, line=AZURE, size=13, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0.55, 3.1, 6, 0.4, "Fault surfaces, one per scenario", size=15, color=NAVY, bold=True)
rows = [
    ("Network", "NSG DENY rule (IaC param)"),
    ("Edge", "App Gateway health probe (IaC param)"),
    ("Compute", "Replica bounds → zero (IaC param)"),
    ("App code", "Slow release: real source → real image"),
    ("Database schema", "Index corruption → REINDEX (live DB drift)"),
    ("Database role", "Connection limit (live DB drift)"),
    ("Secrets", "Rotated Key Vault credential (live drift)"),
    ("Reporting VM", "Grade-export disk fills up (live VM drift)"),
]
y = 3.6
for tier, what in rows:
    add_box(s, 0.55, y, 2.3, 0.42, tier, fill=AZURE, line=None, size=11, bold=True, color=WHITE)
    add_box(s, 2.95, y, 3.6, 0.42, what, fill=WHITE, line=AZURE, size=10, color=INK, align=PP_ALIGN.LEFT)
    y += 0.46
# legend
add_text(s, 7.1, 3.1, 5, 0.4, "Legend", size=15, color=NAVY, bold=True)
legend_chip(s, 7.15, 3.7, AZURE, "Azure service / healthy component")
legend_chip(s, 7.15, 4.1, CYAN, "Container Apps lane")
legend_chip(s, 7.15, 4.5, RED, "Broken component / blocked path")
legend_chip(s, 7.15, 4.9, GREEN, "Good-state flow / reset")
legend_chip(s, 7.15, 5.3, NAVY, "Operations & SRE Agent")
add_box(s, 7.1, 5.95, 5.6, 1.05,
        "Symptom-only by design: alerts and student-facing symptoms never name the cause —\n"
        "diagnosing it is the SRE Agent's job.",
        fill=NAVY, line=None, size=11.5, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

out = Path(os.environ["ZAVA_DECK_OUT"]) if os.environ.get("ZAVA_DECK_OUT") \
    else Path(__file__).resolve().parent / "Zava-Learning-Architecture.pptx"
if not out.is_absolute():
    out = Path(__file__).resolve().parent / out
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
